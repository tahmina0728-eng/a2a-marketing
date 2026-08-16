"""Parser: .pptx (python-pptx)"""
from __future__ import annotations
import base64
import io
import re


def _clean(t: str) -> str:
    return re.sub(r"\s+", " ", t).strip()


def parse_pptx(content: bytes) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError(
            "python-pptx is required for .pptx files. Install: pip install python-pptx"
        )

    prs    = Presentation(io.BytesIO(content))
    blocks: list[dict] = []
    images: list[dict] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                ph_idx = None
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_idx = shape.placeholder_format.idx
                for i, para in enumerate(shape.text_frame.paragraphs):
                    text = _clean(para.text)
                    if not text:
                        continue
                    if ph_idx == 0 and i == 0:
                        blocks.append({"type": "heading", "level": 2, "text": text})
                    else:
                        blocks.append({"type": "paragraph", "text": text})
            if hasattr(shape, "image"):
                try:
                    img  = shape.image
                    b64  = base64.b64encode(img.blob).decode()
                    mime = img.content_type or "image/png"
                    images.append({"b64": b64, "mime": mime})
                except Exception:
                    pass

    return {"blocks": blocks, "images": images[:6]}
