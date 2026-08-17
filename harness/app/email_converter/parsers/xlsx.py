def parse(
    raw: bytes,
    filename: str,
    **kwargs,
) -> dict:

    import io

    try:

        from pptx import Presentation

    except ImportError as exc:

        raise RuntimeError(
            "Install python-pptx"
        ) from exc

    presentation = Presentation(
        io.BytesIO(raw)
    )

    blocks = []

    for slide in (
        presentation.slides
    ):

        for shape in (
            slide.shapes
        ):

            text = getattr(
                shape,
                "text",
                "",
            ).strip()

            if text:

                blocks.append({

                    "type":
                        "paragraph",

                    "text":
                        text,
                })

    return {

        "blocks":
            blocks,

        "images":
            [],

        "image_context":
            [],
    }