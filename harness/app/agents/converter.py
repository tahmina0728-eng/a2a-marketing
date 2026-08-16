"""
converter.py — LEGACY shim. The email converter has been refactored
into app/email_converter/. This file is kept for backward compatibility.

New code should import from app.email_converter directly:
    from app.email_converter import EmailConverterAgent
"""
from __future__ import annotations

import base64
import csv as _csv_mod
import io
import re
from typing import Any


# ── Text utilities ─────────────────────────────────────────────────────────────

def _clean(t: str) -> str:
    return re.sub(r"\s+", " ", t).strip()


def _detect_slot(text: str) -> str | None:
    lower = text.lower().strip()
    for kw in ("subject:", "email subject:", "subject line:"):
        if lower.startswith(kw): return "subject"
    for kw in ("preheader:", "preview text:", "preview:"):
        if lower.startswith(kw): return "preheader"
    for kw in ("headline:", "header:", "title:", "h1:"):
        if lower.startswith(kw): return "headline"
    for kw in ("subline:", "subtitle:", "subheading:", "sub-headline:", "tagline:"):
        if lower.startswith(kw): return "subline"
    for kw in ("cta:", "call to action:", "button:", "action:"):
        if lower.startswith(kw): return "cta"
    return None


def _strip_label(text: str) -> str:
    idx = text.find(":")
    return text[idx + 1:].strip() if idx != -1 else text.strip()


def _escape(text: str) -> str:
    return (
        text
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return r, g, b


def _darken(hex_color: str, factor: float = 0.80) -> str:
    r, g, b = _hex_to_rgb(hex_color)
    return "#{:02x}{:02x}{:02x}".format(
        int(r * factor), int(g * factor), int(b * factor)
    )


def _text_on(hex_color: str) -> str:
    """Return white or near-black depending on luminance."""
    r, g, b = _hex_to_rgb(hex_color)
    lum = 0.299 * r + 0.587 * g + 0.114 * b
    return "#ffffff" if lum < 140 else "#0d0d0d"


# ── Format parsers ─────────────────────────────────────────────────────────────

def _parse_docx(content: bytes) -> dict:
    from docx import Document

    doc = Document(io.BytesIO(content))
    blocks: list[dict] = []
    images: list[dict] = []

    for para in doc.paragraphs:
        text = _clean(para.text)
        if not text:
            continue
        style = (para.style.name or "").lower()
        level = 0
        if "heading 1" in style or "title" in style:
            level = 1
        elif "heading 2" in style:
            level = 2
        elif "heading 3" in style:
            level = 3
        blocks.append({"type": "heading" if level else "paragraph", "level": level, "text": text})

    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                blob = rel.target_part.blob
                mime = rel.target_part.content_type or "image/png"
                images.append({"b64": base64.b64encode(blob).decode(), "mime": mime})
            except Exception:
                pass

    return {"blocks": blocks, "images": images[:6]}


def _parse_pdf(content: bytes) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=content, filetype="pdf")
    blocks: list[dict] = []
    images: list[dict] = []

    for page in doc:
        page_dict = page.get_text("dict")
        for block in page_dict.get("blocks", []):
            btype = block.get("type")
            if btype == 0:  # text
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    line_text = _clean(" ".join(s.get("text", "") for s in spans))
                    if not line_text:
                        continue
                    max_size = max((s.get("size", 0) for s in spans), default=0)
                    if max_size >= 18:
                        blocks.append({"type": "heading", "level": 1, "text": line_text})
                    elif max_size >= 14:
                        blocks.append({"type": "heading", "level": 2, "text": line_text})
                    else:
                        blocks.append({"type": "paragraph", "text": line_text})
            elif btype == 1:  # image
                try:
                    xref = block.get("xref", 0)
                    if xref:
                        img_data = doc.extract_image(xref)
                        raw  = img_data.get("image", b"")
                        mime = f"image/{img_data.get('ext', 'jpeg')}"
                    else:
                        raw  = block.get("image", b"")
                        mime = "image/jpeg"
                    if raw:
                        images.append({"b64": base64.b64encode(raw).decode(), "mime": mime})
                except Exception:
                    pass

    doc.close()
    return {"blocks": blocks, "images": images[:6]}


def _parse_xlsx(content: bytes) -> dict:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    blocks: list[dict] = []
    multi = len(wb.worksheets) > 1

    for sheet in wb.worksheets:
        rows_data: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            cells = [c for c in cells if c]
            if cells:
                rows_data.append(cells)
        if rows_data:
            if multi:
                blocks.append({"type": "heading", "level": 2, "text": sheet.title})
            blocks.append({"type": "table", "headers": rows_data[0], "rows": rows_data[1:]})

    wb.close()
    return {"blocks": blocks, "images": []}


def _parse_csv(content: bytes) -> dict:
    text = content.decode("utf-8-sig", errors="replace")
    reader = _csv_mod.reader(io.StringIO(text))
    rows = [r for r in reader if any(c.strip() for c in r)]
    if not rows:
        return {"blocks": [], "images": []}
    return {"blocks": [{"type": "table", "headers": rows[0], "rows": rows[1:]}], "images": []}


def _parse_txt(content: bytes) -> dict:
    """Plain text: split on blank lines into paragraphs."""
    text = content.decode("utf-8-sig", errors="replace")
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    blocks: list[dict] = []
    for i, para in enumerate(paras):
        lines = [l.strip() for l in para.splitlines() if l.strip()]
        if not lines:
            continue
        first = lines[0]
        rest  = " ".join(lines[1:])
        if i == 0 and len(first) < 120:
            blocks.append({"type": "heading", "level": 1, "text": first})
            if rest:
                blocks.append({"type": "paragraph", "text": rest})
        else:
            blocks.append({"type": "paragraph", "text": " ".join(lines)})
    return {"blocks": blocks, "images": []}


def _parse_image(content: bytes, filename: str, mime: str) -> dict:
    """Treat the uploaded image file itself as a content image."""
    b64  = base64.b64encode(content).decode()
    name = re.sub(r"[-_]+", " ", filename.rsplit(".", 1)[0]).strip().title()
    return {
        "blocks": [{"type": "heading", "level": 1, "text": name}],
        "images": [{"b64": b64, "mime": mime}],
    }


def _parse_pptx(content: bytes) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError(
            "python-pptx is required for .pptx files. "
            "Install it: pip install python-pptx"
        )
    prs = Presentation(io.BytesIO(content))
    blocks: list[dict] = []
    images: list[dict] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                ph_idx = None
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    ph_idx = shape.placeholder_format.idx
                for i, para in enumerate(shape.text_frame.paragraphs):
                    text = _clean(para.text)
                    if not text:
                        continue
                    if ph_idx == 0 and i == 0:
                        blocks.append({"type": "heading", "level": 2, "text": text})
                    else:
                        blocks.append({"type": "paragraph", "text": text})
            if hasattr(shape, "image"):
                try:
                    img  = shape.image
                    b64  = base64.b64encode(img.blob).decode()
                    mime = img.content_type or "image/png"
                    images.append({"b64": b64, "mime": mime})
                except Exception:
                    pass

    return {"blocks": blocks, "images": images[:6]}


# ── Slot mapping ───────────────────────────────────────────────────────────────

def _map_slots(parsed: dict) -> dict[str, Any]:
    slots: dict[str, Any] = {
        "subject":   "",
        "preheader": "",
        "headline":  "",
        "subline":   "",
        "body":      [],
        "cta":       "",
        "tables":    [],
        "images":    parsed.get("images", []),
    }

    for block in parsed.get("blocks", []):
        btype = block.get("type")

        if btype == "table":
            slots["tables"].append({"headers": block.get("headers", []), "rows": block.get("rows", [])})
            continue

        text = block.get("text", "").strip()
        if not text:
            continue

        slot = _detect_slot(text)
        if slot:
            slots[slot] = _strip_label(text)
            continue

        level = block.get("level", 0)
        if btype == "heading" and level == 1:
            if not slots["headline"]:
                slots["headline"] = text
            else:
                slots["body"].append(text)
        elif btype == "heading":
            if slots["headline"] and not slots["subline"] and len(text) < 160:
                slots["subline"] = text
            else:
                slots["body"].append(text)
        else:
            slots["body"].append(text)

    if not slots["subject"] and slots["headline"]:
        slots["subject"] = slots["headline"]

    return slots


def _merge_slots_list(slots_list: list[dict], filenames: list[str]) -> dict[str, Any]:
    """Merge multiple per-file slot dicts into one combined set."""
    if not slots_list:
        return _map_slots({"blocks": [], "images": []})
    if len(slots_list) == 1:
        return slots_list[0]

    first  = slots_list[0]
    merged: dict[str, Any] = {
        "subject":   first.get("subject", ""),
        "preheader": first.get("preheader", ""),
        "headline":  first.get("headline", ""),
        "subline":   first.get("subline", ""),
        "cta":       first.get("cta", ""),
        "body":      list(first.get("body", [])),
        "tables":    list(first.get("tables", [])),
        "images":    list(first.get("images", [])),
        "_sections": [],
    }

    # First file as section 0
    merged["_sections"].append({
        "label":  filenames[0] if filenames else "",
        "body":   list(first.get("body", [])),
        "tables": list(first.get("tables", [])),
        "images": list(first.get("images", [])),
    })

    for slots, fname in zip(slots_list[1:], filenames[1:]):
        merged["_sections"].append({
            "label":  fname,
            "body":   list(slots.get("body", [])),
            "tables": list(slots.get("tables", [])),
            "images": list(slots.get("images", [])),
        })
        merged["body"].extend(slots.get("body", []))
        merged["tables"].extend(slots.get("tables", []))
        merged["images"].extend(slots.get("images", []))
        if not merged["cta"] and slots.get("cta"):
            merged["cta"] = slots["cta"]
        if not merged["subject"] and slots.get("subject"):
            merged["subject"] = slots["subject"]
        if not merged["headline"] and slots.get("headline"):
            merged["headline"] = slots["headline"]

    if not merged["subject"] and merged["headline"]:
        merged["subject"] = merged["headline"]

    merged["images"] = merged["images"][:10]
    return merged


# ── HTML email renderer ────────────────────────────────────────────────────────

def _render_email(
    slots: dict,
    brand_name: str  = "",
    brand_color: str = "#0055A4",
    multi_file: bool = False,
) -> str:
    subject   = _escape(slots.get("subject", ""))
    preheader = slots.get("preheader", "")
    headline  = _escape(slots.get("headline", ""))
    subline   = _escape(slots.get("subline", ""))
    cta       = _escape(slots.get("cta", ""))
    images    = slots.get("images", [])
    sections  = slots.get("_sections", [])
    on_brand  = _text_on(brand_color)
    dark_col  = _darken(brand_color, 0.82)
    body_raw  = slots.get("body", [])

    # ── Preheader (hidden preview text for inbox) ──
    _ph = preheader or (body_raw[0][:120] if body_raw else "") or slots.get("headline", "")
    preheader_html = (
        f'<div style="display:none;max-height:0;overflow:hidden;mso-hide:all;'
        f'font-size:1px;color:#f5f5f5;line-height:1px;">'
        f'{_escape(_ph)}&nbsp;' + '&#847;&nbsp;' * 80 + '</div>'
        if _ph else ""
    )

    # ── Classify body into paragraphs / bullets / subheadings ──
    # Items < 90 chars → bullet point; items ending ":" → subheading; rest → paragraph
    classified: list[tuple[str, str]] = []
    for item in body_raw:
        s = item.strip()
        if len(s) < 4:
            continue
        if s.endswith(":") and len(s) < 80:
            classified.append(("subhead", s.rstrip(":")))
        elif len(s) < 90:
            classified.append(("bullet", s))
        else:
            classified.append(("para", s))

    # Group consecutive bullets together
    groups: list[tuple[str, list[str]]] = []
    for (itype, itext) in classified:
        if groups and groups[-1][0] == "bullet" and itype == "bullet":
            groups[-1][1].append(itext)
        else:
            groups.append((itype, [itext]))

    # ── Shared inner helpers ──
    brand_label = _escape(brand_name) if brand_name else "CampaignOS"

    def _eyebrow(label: str) -> str:
        return (
            f'<tr><td style="padding:28px 40px 10px;" class="pad">'
            f'<p style="margin:0;font-size:11px;font-weight:800;letter-spacing:.11em;'
            f'text-transform:uppercase;color:{brand_color};'
            f'font-family:Helvetica,Arial,sans-serif;">{_escape(label)}</p>'
            f'</td></tr>'
        )

    def _para_block(text: str, lead: bool = False) -> str:
        sz  = "17px" if lead else "15.5px"
        fw  = "500"  if lead else "400"
        col = "#333333" if lead else "#454545"
        return (
            f'<tr><td style="padding:0 40px 16px;" class="pad">'
            f'<p style="margin:0;font-size:{sz};line-height:1.8;color:{col};font-weight:{fw};'
            f'font-family:\'Helvetica Neue\',Helvetica,Arial,sans-serif;">'
            f'{_escape(text)}</p></td></tr>'
        )

    def _subhead_block(text: str) -> str:
        return (
            f'<tr><td style="padding:20px 40px 6px;" class="pad">'
            f'<p style="margin:0;font-size:13px;font-weight:800;letter-spacing:.05em;'
            f'text-transform:uppercase;color:{brand_color};'
            f'font-family:\'Helvetica Neue\',Helvetica,Arial,sans-serif;">'
            f'{_escape(text)}</p></td></tr>'
        )

    def _bullet_list(items: list[str]) -> str:
        rows = ""
        for item in items:
            rows += (
                f'<tr>'
                f'<td width="28" valign="top" style="padding-bottom:13px;padding-top:2px;">'
                f'<table cellpadding="0" cellspacing="0"><tr>'
                f'<td style="width:8px;height:8px;background:{brand_color};'
                f'border-radius:50%;font-size:0;line-height:0;">&nbsp;</td>'
                f'</tr></table></td>'
                f'<td style="padding-bottom:13px;font-size:15px;line-height:1.7;'
                f'color:#333333;font-family:\'Helvetica Neue\',Helvetica,Arial,sans-serif;">'
                f'{_escape(item)}</td></tr>'
            )
        return (
            f'<tr><td style="padding:4px 40px 4px;" class="pad">'
            f'<table width="100%" cellpadding="0" cellspacing="0">{rows}</table>'
            f'</td></tr>'
        )

    def _table_block(tbl: dict) -> str:
        headers = tbl.get("headers", [])
        rows    = tbl.get("rows", [])
        if not headers and not rows:
            return ""
        th_cells = "".join(
            f'<th style="background:{brand_color};color:{on_brand};padding:11px 14px;'
            f'text-align:left;font-size:12px;font-weight:700;white-space:nowrap;'
            f'font-family:Helvetica,Arial,sans-serif;">{_escape(str(h))}</th>'
            for h in headers
        )
        tbody_rows = ""
        for i, row in enumerate(rows[:50]):
            bg  = "#f7f8fc" if i % 2 == 0 else "#ffffff"
            tds = "".join(
                f'<td style="padding:10px 14px;font-size:13px;color:#444444;'
                f'border-bottom:1px solid #eeeeee;'
                f'font-family:Helvetica,Arial,sans-serif;">{_escape(str(c))}</td>'
                for c in row
            )
            tbody_rows += f'<tr style="background:{bg};">{tds}</tr>'
        note = (
            f'<p style="margin:6px 0 0;font-size:11px;color:#999;'
            f'font-family:Helvetica,Arial,sans-serif;">Showing 50 of {len(rows)} rows.</p>'
            if len(rows) > 50 else ""
        )
        return (
            f'<tr><td style="padding:12px 40px 4px;" class="pad">'
            f'<div style="overflow-x:auto;-webkit-overflow-scrolling:touch;">'
            f'<table width="100%" cellpadding="0" cellspacing="0" '
            f'style="border-collapse:collapse;border:1px solid #e8e8e8;min-width:300px;">'
            f'<thead><tr>{th_cells}</tr></thead><tbody>{tbody_rows}</tbody></table>'
            f'</div>{note}</td></tr>'
        )

    def _section_divider(label: str) -> str:
        return (
            f'<tr><td style="padding:28px 40px 0;" class="pad">'
            f'<table width="100%" cellpadding="0" cellspacing="0"><tr>'
            f'<td style="border-top:1px solid #e8e8e8;font-size:0;">&nbsp;</td>'
            f'</tr></table></td></tr>'
            f'<tr><td style="padding:16px 40px 4px;" class="pad">'
            f'<table cellpadding="0" cellspacing="0"><tr>'
            f'<td style="background:{brand_color};width:3px;border-radius:3px;">&nbsp;</td>'
            f'<td style="padding-left:12px;font-size:12px;font-weight:800;'
            f'letter-spacing:.08em;text-transform:uppercase;color:#777777;'
            f'font-family:Helvetica,Arial,sans-serif;">{_escape(label)}</td>'
            f'</tr></table></td></tr>'
        )

    # ── Brand header ──
    brand_html = f"""
        <tr>
          <td style="background:{brand_color};">
            <table width="100%" cellpadding="0" cellspacing="0">
              <tr>
                <td style="padding:18px 40px;">
                  <table width="100%" cellpadding="0" cellspacing="0">
                    <tr>
                      <td>
                        <p style="margin:0;font-size:17px;font-weight:900;
                           letter-spacing:.09em;text-transform:uppercase;color:{on_brand};
                           font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
                          {brand_label}
                        </p>
                      </td>
                      <td align="right">
                        <p style="margin:0;font-size:10px;color:{on_brand};opacity:.65;
                           font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;
                           letter-spacing:.07em;text-transform:uppercase;">
                          Campaign&nbsp;Update
                        </p>
                      </td>
                    </tr>
                  </table>
                </td>
              </tr>
              <tr>
                <td style="height:3px;background:{dark_col};font-size:0;">&nbsp;</td>
              </tr>
            </table>
          </td>
        </tr>"""

    # ── Hero ──
    # If we have an image: show it full-width, then headline below on white bg.
    # If no image: brand-color gradient banner with headline text in white (more visual impact).
    gallery_images = images[:]
    hl_rendered_in_hero = False

    if gallery_images:
        hero = gallery_images.pop(0)
        hero_html = f"""
        <tr>
          <td style="padding:0;line-height:0;font-size:0;background:{brand_color};">
            <img src="data:{hero['mime']};base64,{hero['b64']}"
                 width="600" alt="{brand_label}"
                 style="display:block;width:100%;max-width:600px;height:auto;border:0;" />
          </td>
        </tr>"""
        # Headline accent strip immediately under the image — ties them together
        if headline:
            hero_html += f"""
        <tr>
          <td style="background:{brand_color};padding:28px 40px 26px;" class="pad">
            <h1 style="margin:0 0 8px;font-size:32px;font-weight:900;line-height:1.2;
               letter-spacing:-.02em;color:{on_brand};
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
              {headline}
            </h1>"""
            if subline:
                hero_html += f"""
            <p style="margin:0;font-size:17px;line-height:1.55;color:{on_brand};opacity:.85;
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
              {subline}
            </p>"""
            hero_html += """
          </td>
        </tr>"""
            hl_rendered_in_hero = True
    else:
        # Gradient banner — headline in white, no image
        hl_in_banner = ""
        if headline:
            hl_in_banner += (
                f'<h1 style="margin:0 0 10px;font-size:36px;font-weight:900;line-height:1.18;'
                f'letter-spacing:-.025em;color:{on_brand};'
                f'font-family:\'Helvetica Neue\',Helvetica,Arial,sans-serif;">{headline}</h1>'
            )
        if subline:
            hl_in_banner += (
                f'<p style="margin:0;font-size:18px;line-height:1.55;color:{on_brand};opacity:.82;'
                f'font-family:\'Helvetica Neue\',Helvetica,Arial,sans-serif;">{subline}</p>'
            )
        hero_html = f"""
        <tr>
          <td style="background:linear-gradient(135deg,{brand_color} 0%,{dark_col} 100%);
                     padding:52px 40px 44px;" class="pad">
            {hl_in_banner}
          </td>
        </tr>"""
        hl_rendered_in_hero = True

    # Standalone headline block (only when image was present and headline NOT in hero strip)
    hl_html = ""
    if headline and not hl_rendered_in_hero:
        hl_html = f"""
        <tr>
          <td style="padding:40px 40px 0;" class="pad">
            <h1 style="margin:0 0 12px;font-size:34px;font-weight:900;line-height:1.2;
               letter-spacing:-.025em;color:#0d0d0d;
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">{headline}</h1>
            <div style="width:44px;height:4px;background:{brand_color};border-radius:99px;"></div>
          </td>
        </tr>"""
        if subline:
            hl_html += f"""
        <tr>
          <td style="padding:14px 40px 0;" class="pad">
            <p style="margin:0;font-size:18px;line-height:1.55;color:#555555;
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">{subline}</p>
          </td>
        </tr>"""

    spacer_after_hero = '<tr><td style="height:28px;font-size:0;">&nbsp;</td></tr>'

    # ── Build body HTML ──
    body_html   = ""
    tables_html = ""

    if multi_file and sections:
        for i, sec in enumerate(sections):
            if i > 0 and sec.get("label"):
                body_html += _section_divider(sec["label"])

            # Classify this section's body
            sec_classified: list[tuple[str, str]] = []
            for item in sec.get("body", []):
                s = item.strip()
                if len(s) < 4: continue
                if s.endswith(":") and len(s) < 80:
                    sec_classified.append(("subhead", s.rstrip(":")))
                elif len(s) < 90:
                    sec_classified.append(("bullet", s))
                else:
                    sec_classified.append(("para", s))

            sec_groups: list[tuple[str, list[str]]] = []
            for (itype, itext) in sec_classified:
                if sec_groups and sec_groups[-1][0] == "bullet" and itype == "bullet":
                    sec_groups[-1][1].append(itext)
                else:
                    sec_groups.append((itype, [itext]))

            lead        = True
            bullet_done = False
            for (gtype, gitems) in sec_groups:
                if gtype == "para":
                    body_html += _para_block(gitems[0], lead=(lead and i == 0))
                    lead = False
                elif gtype == "subhead":
                    body_html += _subhead_block(gitems[0])
                elif gtype == "bullet":
                    if not bullet_done:
                        body_html += _eyebrow("Key Points")
                        bullet_done = True
                    body_html += _bullet_list(gitems)

            for tbl in sec.get("tables", []):
                tables_html += _table_block(tbl)

            # Extra images in this section (skip the overall hero from section 0)
            sec_imgs = sec.get("images", [])
            skip = 1 if (i == 0 and images) else 0
            for img in sec_imgs[skip:]:
                body_html += (
                    f'<tr><td style="padding:16px 40px 0;" class="pad">'
                    f'<img src="data:{img["mime"]};base64,{img["b64"]}" alt="" width="520"'
                    f' style="display:block;width:100%;max-width:520px;height:auto;'
                    f'border-radius:6px;border:1px solid #eeeeee;" /></td></tr>'
                )
    else:
        lead        = True
        bullet_done = False
        for (gtype, gitems) in groups:
            if gtype == "para":
                body_html += _para_block(gitems[0], lead=lead)
                lead = False
            elif gtype == "subhead":
                body_html += _subhead_block(gitems[0])
            elif gtype == "bullet":
                if not bullet_done:
                    body_html += _eyebrow("Key Points")
                    bullet_done = True
                body_html += _bullet_list(gitems)

        src_tables = slots.get("tables", [])
        if src_tables:
            tables_html += _eyebrow("Data")
            for tbl in src_tables:
                tables_html += _table_block(tbl)

    # ── Gallery (remaining images, 2-col grid) ──
    gallery_html = ""
    if gallery_images:
        pairs = [gallery_images[i:i+2] for i in range(0, len(gallery_images), 2)]
        for pair in pairs:
            left = pair[0]
            right_td = '<td width="51%">&nbsp;</td>'
            if len(pair) > 1:
                r = pair[1]
                right_td = (
                    f'<td width="2%" style="font-size:0;">&nbsp;</td>'
                    f'<td width="49%" class="col2" valign="top">'
                    f'<img src="data:{r["mime"]};base64,{r["b64"]}" alt="" width="260"'
                    f' style="display:block;width:100%;max-width:260px;height:auto;'
                    f'border-radius:6px;border:1px solid #eeeeee;" /></td>'
                )
            gallery_html += (
                f'<tr><td style="padding:16px 40px 0;" class="pad">'
                f'<table width="100%" cellpadding="0" cellspacing="0"><tr>'
                f'<td width="49%" class="col2" valign="top">'
                f'<img src="data:{left["mime"]};base64,{left["b64"]}" alt="" width="260"'
                f' style="display:block;width:100%;max-width:260px;height:auto;'
                f'border-radius:6px;border:1px solid #eeeeee;" /></td>'
                f'{right_td}</tr></table></td></tr>'
            )

    # ── CTA button ──
    if cta:
        cta_html = f"""
        <tr>
          <td align="center" style="padding:36px 40px 12px;">
            <!--[if mso]>
            <v:roundrect xmlns:v="urn:schemas-microsoft-com:vml"
              href="#" style="height:52px;v-text-anchor:middle;width:240px;"
              arcsize="10%" stroke="f" fillcolor="{brand_color}">
            <w:anchorlock/><center><![endif]-->
            <a href="#"
               style="display:inline-block;background:{brand_color};color:{on_brand};
                      font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;
                      font-size:16px;font-weight:800;text-decoration:none;
                      padding:16px 50px;border-radius:8px;letter-spacing:.03em;">
              {cta}&nbsp;&rarr;
            </a>
            <!--[if mso]></center></v:roundrect><![endif]-->
          </td>
        </tr>
        <tr><td style="height:16px;font-size:0;">&nbsp;</td></tr>"""
    else:
        cta_html = '<tr><td style="height:36px;font-size:0;">&nbsp;</td></tr>'

    # ── Footer ──
    footer_brand = _escape(brand_name) if brand_name else "CampaignOS"
    footer_html = f"""
        <tr>
          <td style="border-top:1px solid #e8e8e8;padding:28px 40px 36px;">
            <p style="margin:0 0 10px;font-size:12px;color:#aaaaaa;text-align:center;
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;line-height:1.7;">
              You are receiving this because you subscribed to {footer_brand} updates.
            </p>
            <p style="margin:0 0 12px;font-size:12px;color:#aaaaaa;text-align:center;
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
              <a href="#" style="color:#aaaaaa;text-decoration:underline;">Unsubscribe</a>
              &nbsp;&middot;&nbsp;
              <a href="#" style="color:#aaaaaa;text-decoration:underline;">Privacy&nbsp;Policy</a>
              &nbsp;&middot;&nbsp;
              <a href="#" style="color:#aaaaaa;text-decoration:underline;">View&nbsp;in&nbsp;browser</a>
            </p>
            <p style="margin:0;font-size:11px;color:#cccccc;text-align:center;
               font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;">
              &copy; 2026 {footer_brand}. Powered by CampaignOS.
            </p>
          </td>
        </tr>"""

    return f"""<!DOCTYPE html>
<html lang="en" xmlns:v="urn:schemas-microsoft-com:vml" xmlns:o="urn:schemas-microsoft-com:office:office">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width,initial-scale=1.0">
  <meta name="x-apple-disable-message-reformatting">
  <meta http-equiv="X-UA-Compatible" content="IE=edge">
  <title>{subject}</title>
  <!--[if mso]>
  <noscript><xml><o:OfficeDocumentSettings>
    <o:PixelsPerInch>96</o:PixelsPerInch>
  </o:OfficeDocumentSettings></xml></noscript>
  <![endif]-->
  <style>
    body{{margin:0;padding:0;background:#ebebeb;-webkit-text-size-adjust:100%;}}
    img{{border:0;outline:none;text-decoration:none;display:block;}}
    table{{border-collapse:collapse;mso-table-lspace:0pt;mso-table-rspace:0pt;}}
    a{{color:{brand_color};}}
    @media only screen and (max-width:640px){{
      .wrap{{width:100%!important;max-width:100%!important;}}
      .pad{{padding-left:24px!important;padding-right:24px!important;}}
      h1{{font-size:26px!important;line-height:1.25!important;}}
      .col2{{display:block!important;width:100%!important;max-width:100%!important;margin-bottom:12px!important;}}
    }}
  </style>
</head>
<body style="margin:0;padding:0;background:#ebebeb;">
{preheader_html}
<table width="100%" cellpadding="0" cellspacing="0" style="background:#ebebeb;min-width:280px;" role="presentation">
  <tr>
    <td align="center" style="padding:28px 16px 48px;">
      <table class="wrap" width="600" cellpadding="0" cellspacing="0"
        style="background:#ffffff;border-radius:10px;overflow:hidden;
               box-shadow:0 4px 32px rgba(0,0,0,.14);" role="presentation">
        {brand_html}
        {hero_html}
        {hl_html}
        {spacer_after_hero}
        {body_html}
        {tables_html}
        {gallery_html}
        {cta_html}
        {footer_html}
      </table>
    </td>
  </tr>
</table>
</body>
</html>"""


# ── Parser registry ────────────────────────────────────────────────────────────

_PARSERS: dict[str, Any] = {
    "docx": _parse_docx,
    "doc":  _parse_docx,     # try python-docx; will fail on true binary .doc
    "pdf":  _parse_pdf,
    "xlsx": _parse_xlsx,
    "xls":  _parse_xlsx,
    "csv":  _parse_csv,
    "txt":  _parse_txt,
    "pptx": _parse_pptx,
    "jpg":  None,            # handled via _parse_image
    "jpeg": None,
    "png":  None,
    "gif":  None,
    "webp": None,
}

_IMAGE_EXTS = {"jpg", "jpeg", "png", "gif", "webp"}

_EXT_MIME: dict[str, str] = {
    "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "png": "image/png",  "gif": "image/gif",
    "webp": "image/webp",
}


def _parse_one(content: bytes, filename: str) -> dict:
    """Parse a single file. Returns {blocks, images}."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in _IMAGE_EXTS:
        mime = _EXT_MIME.get(ext, "image/jpeg")
        return _parse_image(content, filename, mime)
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(
            f"Unsupported file type: .{ext}  "
            f"(supported: docx, pdf, xlsx, xls, csv, txt, pptx, jpg, jpeg, png, gif, webp)"
        )
    return parser(content)


# ── Public entry points ────────────────────────────────────────────────────────

def convert_documents(
    files: list[tuple[bytes, str]],
    brand_name: str  = "",
    brand_color: str = "#0055A4",
) -> dict:
    """
    Convert one or more (content_bytes, filename) pairs into a single HTML email.

    Returns:
      { html, slots, image_count, filename, file_count }
    """
    if not files:
        raise ValueError("No files provided.")

    slots_list: list[dict] = []
    filenames: list[str]   = []

    for content, filename in files:
        try:
            parsed = _parse_one(content, filename)
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError(f"Could not parse {filename}: {exc}") from exc
        slots_list.append(_map_slots(parsed))
        filenames.append(filename)

    multi = len(slots_list) > 1
    merged = _merge_slots_list(slots_list, filenames) if multi else slots_list[0]

    html = _render_email(
        merged,
        brand_name=brand_name,
        brand_color=brand_color,
        multi_file=multi,
    )

    # Public slots (strip internal _sections and raw image bytes)
    public_slots = {
        k: v for k, v in merged.items()
        if k not in ("images", "_sections")
    }

    combined_name = (
        filenames[0] if not multi
        else f"Combined ({len(filenames)} files)"
    )

    return {
        "html":        html,
        "slots":       public_slots,
        "image_count": len(merged.get("images", [])),
        "filename":    combined_name,
        "file_count":  len(files),
    }


def convert_document(
    content: bytes,
    filename: str,
    brand_name: str  = "",
    brand_color: str = "#0055A4",
) -> dict:
    """Single-file convenience wrapper — backward compatible."""
    return convert_documents([(content, filename)], brand_name=brand_name, brand_color=brand_color)
